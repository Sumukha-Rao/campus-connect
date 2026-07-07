# -*- coding: utf-8 -*-
"""Fill the SEE_LAB_MAD.pptx template with College Forum App content + speaker
notes + screenshots. Template design (index/header/footer/background/graphics)
is left untouched; only placeholder text, notes and a few images are added."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from PIL import Image

TPL   = r"C:\Users\Lenovo\Downloads\SEE_LAB_MAD.pptx"
OUT   = r"D:\MCA\MAD\Project\campus-connect\College_Forum_App_Presentation.pptx"
SHOTS = r"D:\MCA\MAD\Project\campus-connect\report_assets\shots"
FONT  = "Bookman Old Style"

prs = Presentation(TPL)
SW, SH = prs.slide_width, prs.slide_height

def body_ph(slide):
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.idx == 1:
            return sh
    return None

def fill_body(slide, items):
    """items: list of (text, level). Inherits the template's bullet/list style."""
    tf = body_ph(slide).text_frame
    tf.word_wrap = True
    tf.clear()
    first = True
    for text, lvl in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = text
        p.level = lvl
        first = False

def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def place_image(slide, path, x, y, w, h, caption=None):
    """Contain the image inside box (x,y,w,h) in inches, centered."""
    iw, ih = Image.open(path).size
    bw, bh = Inches(w), Inches(h)
    scale = min(bw / iw, bh / ih)
    dw, dh = int(iw * scale), int(ih * scale)
    px = Inches(x) + (bw - dw) // 2
    py = Inches(y) + (bh - dh) // 2
    slide.shapes.add_picture(path, px, py, dw, dh)
    if caption:
        tb = slide.shapes.add_textbox(Inches(x), Inches(y) + bh, bw, Inches(0.3))
        tfx = tb.text_frame; tfx.word_wrap = True
        para = tfx.paragraphs[0]; para.alignment = 2  # center
        r = para.add_run(); r.text = caption
        r.font.size = Pt(11); r.font.name = FONT; r.font.italic = True

S = prs.slides

# ---------------- Slide 0 : Title ----------------
title = None
for sh in S[0].shapes:
    if sh.is_placeholder and sh.placeholder_format.idx == 0:
        title = sh
# replace just the 3rd run ("Title of Project") -> project title
runs = title.text_frame.paragraphs[0].runs
runs[-1].text = "College Forum App"

sub = None
for sh in S[0].shapes:
    if sh.is_placeholder and sh.placeholder_format.idx == 1:
        sub = sh
stf = sub.text_frame; stf.word_wrap = True; stf.clear()
sub_lines = [
    ("Team Members:", True),
    ("Priya N (1RV25MC075)", False),
    ("Spandana N (1RV25MC096)", False),
    ("Sumanth Bhaskar Hegde (1RV25MC103)", False),
    ("Sumukha Rao B S (1RV25MC104)", False),
    ("Guide: Prof. Prashanth K, Assistant Professor, Dept. of MCA", True),
    ("Submission Date: July 2026", True),
]
first = True
for text, bold in sub_lines:
    p = stf.paragraphs[0] if first else stf.add_paragraph()
    p.alignment = 1  # left
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.bold = bold; r.font.size = Pt(16)
    first = False
set_notes(S[0],
    "Good morning. We are Priya, Spandana, Sumanth and Sumukha. Our Mobile Application "
    "Development project is the College Forum App, a Progressive Web Application that acts "
    "as the official, unified notice board for RV College of Engineering. It is built under "
    "the guidance of Prof. Prashanth K.")

# ---------------- Slide 1 : Index ----------------
fill_body(S[1], [
    ("Introduction", 0),
    ("Key Features", 0),
    ("Technology Stack", 0),
    ("Requirements & Features", 0),
    ("User Interface Mockups / Wireframes", 0),
    ("Demonstration", 0),
    ("Conclusion", 0),
])
set_notes(S[1],
    "This is the flow of our presentation. We begin with the problem and an introduction, "
    "move through the key features, technology stack and requirements, then show the user "
    "interface and a live demonstration, and finish with the conclusion and future scope.")

# ---------------- Slide 2 : Introduction ----------------
fill_body(S[2], [
    ("At RVCE, important updates are scattered across WhatsApp groups, emails and physical boards.", 0),
    ("College Forum App is a Progressive Web Application that centralizes all campus notices.", 0),
    ("In the user interface it is branded as “RVCE Connect”.", 1),
    ("Three roles: Admin, Publisher (faculty / HODs / club heads) and Viewer (student).", 0),
    ("Departments and clubs act as communities that students subscribe to.", 0),
    ("Mobile-first, installable, and works offline.", 0),
])
set_notes(S[2],
    "Communication at RVCE is fragmented: a student in one department often misses a seminar, "
    "hackathon or placement talk organised by another. The College Forum App solves this by "
    "providing one centralized, role-controlled notice board. Only authorised publishers can "
    "broadcast, while students get a clean, personalized, read-only feed. The app is branded "
    "RVCE Connect in its interface, is installable like a native app, and keeps working offline.")

# ---------------- Slide 3 : Key Features ----------------
fill_body(S[3], [
    ("Role-based access control – Admin, Publisher, Viewer.", 0),
    ("Community-targeted notices for departments and clubs.", 0),
    ("Subscriptions, a personalized feed and post likes.", 0),
    ("Real-time Web Push notifications via a per-community bell.", 0),
    ("Post scheduling and automatic expiry / archiving.", 0),
    ("Admin-only post deletion, recorded in an audit log.", 1),
    ("PWA: offline access through service workers; installable.", 0),
    ("Secure JWT authentication with bcrypt-hashed passwords.", 0),
])
set_notes(S[3],
    "The core features map directly to the problem. Role-based access keeps the signal-to-noise "
    "ratio high. Publishers target a specific community – a department or a club – so students "
    "only see what is relevant. Students subscribe, like posts, and can switch on a bell for "
    "genuine browser push notifications using the Web Push protocol. Posts can be scheduled and "
    "auto-expire, deletions are admin-only and audit-logged, and the whole thing is a secure, "
    "offline-capable PWA.")

# ---------------- Slide 4 : Technology Stack ----------------
fill_body(S[4], [
    ("Frontend: HTML5, CSS3, Bootstrap 5, Vanilla JavaScript.", 0),
    ("Backend: Node.js with Express.js (RESTful APIs).", 0),
    ("Database: MySQL 8.", 0),
    ("Authentication & security: JSON Web Tokens (JWT), bcrypt.", 0),
    ("PWA: Web App Manifest, Service Workers, Cache Storage, Web Push (VAPID).", 0),
    ("Tooling: VS Code, npm, Git & GitHub, Docker (MySQL).", 0),
])
set_notes(S[4],
    "We used a modern full-stack web architecture with no heavy SPA framework. The frontend is "
    "plain HTML, CSS, Bootstrap 5 and vanilla JavaScript. The backend is Node.js and Express "
    "exposing REST APIs, backed by MySQL 8. Security uses JWT and bcrypt. The PWA layer – manifest, "
    "service workers, cache storage and VAPID Web Push – is what makes it installable, offline-ready "
    "and able to push notifications. We managed the project with npm, Git and Docker.")

# ---------------- Slide 5 : Requirements & Features ----------------
fill_body(S[5], [
    ("Functional requirements:", 0),
    ("Authentication, registration and role-based access (JWT).", 1),
    ("User and community management by the administrator.", 1),
    ("Notice publishing targeted to chosen communities.", 1),
    ("Personalized feed, likes and subscriptions.", 1),
    ("Push notifications and offline access.", 1),
    ("Non-functional requirements:", 0),
    ("Performance, security, scalability, usability and responsive compatibility.", 1),
])
set_notes(S[5],
    "On the functional side the system must authenticate users with role-based access, let admins "
    "manage users and communities, let publishers post to targeted audiences, and let students "
    "view a personalized feed, like posts and subscribe. Non-functionally it must be fast, secure, "
    "scalable, usable and consistent across desktop, tablet and mobile browsers.")

# ---------------- Slide 6 : UI Mockups / Wireframes ----------------
fill_body(S[6], [
    ("Navigation flow:", 0),
    ("Login authenticates and routes users to a role-based dashboard.", 1),
    ("Viewer: Feed ↔ Communities – subscribe, like and read notices.", 1),
    ("Publisher: compose announcements; Admin: dashboard and moderation.", 1),
])
# resize body to upper-left so screenshots sit below
bp = body_ph(S[6])
bp.top = Inches(2.35); bp.left = Inches(0.5); bp.width = Inches(12.3); bp.height = Inches(1.6)
row_y, row_h = 4.2, 2.9
cells = [
    ("01_login.png", "Login screen"),
    ("03_viewer_feed.png", "Student feed"),
    ("05_publisher_compose.png", "Publisher – compose"),
    ("07_admin_dashboard.png", "Admin dashboard"),
]
x = 0.55
for fname, cap in cells:
    place_image(S[6], SHOTS + "\\" + fname, x, row_y, 3.0, row_h, caption=cap)
    x += 3.18
set_notes(S[6],
    "Here are the actual screens. After login the user is routed by role. A student lands on the "
    "Feed, can switch to Communities to subscribe and turn on the bell, and can like posts. A "
    "publisher gets a Compose screen to create an announcement for a community, with optional "
    "scheduling and expiry. An admin gets a dashboard with the member directory, ban/promote "
    "controls and community management. The interface is mobile-first and theme-aware.")

# ---------------- Slide 7 : Demonstration ----------------
fill_body(S[7], [
    ("Live (development): http://localhost:3000", 0),
    ("Flow: register / login → browse feed → like & subscribe.", 0),
    ("Publisher posts a notice → subscribers get a push notification.", 0),
    ("Disconnect the network → cached notices remain available (offline).", 0),
])
bp7 = body_ph(S[7])
bp7.top = Inches(2.35); bp7.left = Inches(0.5); bp7.width = Inches(12.3); bp7.height = Inches(1.6)
x = 2.0
for fname, cap in (("02_register.png", "Self-registration"), ("06_publisher_feed.png", "Campus feed")):
    place_image(S[7], SHOTS + "\\" + fname, x, 4.2, 4.2, 2.9, caption=cap)
    x += 5.2
set_notes(S[7],
    "For the demo we run the app locally at localhost:3000. We will register a new student, log in, "
    "browse the feed, like a post and subscribe to a community. Then, as a publisher, we will post "
    "a notice and show the push notification reaching a subscriber. Finally we will switch off the "
    "network and refresh to show that previously loaded notices are still available offline.")

# ---------------- Slide 8 : Conclusion ----------------
fill_body(S[8], [
    ("Delivered a secure, centralized campus communication platform.", 0),
    ("Role-based broadcasting reduces information overload.", 0),
    ("PWA adds offline access, push notifications and installability.", 0),
    ("Learning: full-stack development, database design, JWT auth, PWA and UI/UX.", 0),
    ("Future scope: event RSVP, analytics, advanced search, ERP/LMS integration.", 0),
])
set_notes(S[8],
    "In conclusion, the College Forum App successfully delivers a secure and efficient campus "
    "communication system. Role-based broadcasting keeps the feed relevant, and the PWA layer adds "
    "offline access, push notifications and installability. Building it gave us hands-on experience "
    "with full-stack development, relational database design, token authentication and PWA "
    "techniques. In future we plan to add event RSVPs, analytics, advanced search and integration "
    "with institutional ERP and LMS systems.")

# ---------------- Slide 9 : Thank You ----------------
set_notes(S[9], "Thank you for your attention. We are happy to take any questions.")

prs.save(OUT)
print("PPT BUILT ->", OUT, "| slides:", len(prs.slides))
